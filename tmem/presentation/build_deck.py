#!/usr/bin/env python3
"""Build T_MEM_Deck.pptx for thesis / conference presentation."""

from __future__ import annotations

import csv
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
RESULTS_CSV = HERE / "comparison_table.csv"


def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def _two_column_text(prs: Presentation, title: str, left_title: str, left_lines: list[str], right_title: str, right_lines: list[str]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    # Left column
    lx, ly, lw, lh = Inches(0.5), Inches(1.1), Inches(4.4), Inches(5.5)
    left = slide.shapes.add_textbox(lx, ly, lw, lh)
    ltf = left.text_frame
    ltf.word_wrap = True
    p0 = ltf.paragraphs[0]
    p0.text = left_title
    p0.font.size = Pt(20)
    p0.font.bold = True
    for line in left_lines:
        p = ltf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    # Right column
    rx, ry, rw, rh = Inches(5.1), Inches(1.1), Inches(4.4), Inches(5.5)
    right = slide.shapes.add_textbox(rx, ry, rw, rh)
    rtf = right.text_frame
    rtf.word_wrap = True
    p0 = rtf.paragraphs[0]
    p0.text = right_title
    p0.font.size = Pt(20)
    p0.font.bold = True
    for line in right_lines:
        p = rtf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.space_after = Pt(6)


def _equation_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.55))
    box.text_frame.text = title
    box.text_frame.paragraphs[0].font.size = Pt(26)
    box.text_frame.paragraphs[0].font.bold = True
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(8.8), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.name = "Courier New"
        p.space_after = Pt(10)


def _table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.55))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    cols, rcount = len(headers), len(rows) + 1
    tbl = slide.shapes.add_table(rcount, cols, Inches(0.35), Inches(1.0), Inches(9.3), Inches(0.35 * rcount + 0.2)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)[:80]
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)


def _image_slide(prs: Presentation, title: str, image_path: Path, max_h: float = 5.2) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(22)
    tbox.text_frame.paragraphs[0].font.bold = True
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.45), Inches(0.95), height=Inches(max_h))
    else:
        miss = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(8), Inches(1))
        miss.text_frame.text = f"(Missing image: {image_path.name})"


def load_csv_rows() -> list[dict]:
    if not RESULTS_CSV.exists():
        return []
    with open(RESULTS_CSV, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(
        prs,
        "T-MEM: Temporal-aware Memory",
        "Module 1 — Temporal relevance scoring over A-MEM (Xu et al., NeurIPS 2025)\nLoCoMo-style simulation • Self-contained demo (no LLM API)",
    )

    _bullet_slide(
        prs,
        "Motivation",
        [
            "A-MEM retrieval ranks memories by cosine similarity only (paper Eq. 9).",
            "No explicit temporal signal — strong on Temporal LoCoMo category but improvable.",
            "T-MEM adds a decay / usage / link-structure relevance term blended with cosine.",
            "Setting α = 1 recovers A-MEM exactly for baseline comparison.",
        ],
    )

    _two_column_text(
        prs,
        "Memory note (A-MEM Eq. 1 + T-MEM)",
        "Paper fields",
        [
            "c — content",
            "t — timestamp",
            "K — keywords",
            "G — tags",
            "X — contextual sentence",
            "e — embedding of concat(c,K,G,X)",
            "L — linked note IDs",
        ],
        "T-MEM extensions",
        [
            "access_count — retrieval hits",
            "last_accessed — bookkeeping",
            "relevance_score — optional cache",
            "Links simulated: cosine > 0.5 on e (LLM linking proxy)",
        ],
    )

    _equation_slide(
        prs,
        "Temporal relevance and retrieval",
        [
            "r_i = decay(m_i) × reinforce(m_i) × link_bonus(m_i)",
            "",
            "decay     = exp(−λ × age)     [age = month-equivalents: elapsed days / 30]",
            "reinforce = 1 + β × log(1 + access_count)",
            "link_bonus = 1 + γ × (|L_i| / max_j |L_j|)",
            "",
            "score(q, m_i) = α × cos(e_q, e_i) + (1 − α) × r̂_i",
            "r̂_i = min-max normalize r_i over all memories M",
            "",
            "Defaults: λ=0.1, β=0.5, γ=0.3, α=0.7 • Encoder: all-MiniLM-L6-v2",
        ],
    )

    _bullet_slide(
        prs,
        "Simulation (thesis demo)",
        [
            "~30 memory notes across 10 sessions (weeks–months timeline).",
            "Scenarios: job change (Google→Meta), relocation (NYC→Seattle), photography hobby, Japan trip, contradictions.",
            "Photography notes get higher access_count (spacing / reinforcement).",
            "7 evaluation queries across temporal, multi-hop, adversarial, single-hop, open-domain.",
        ],
    )

    rows_data = load_csv_rows()
    table_rows: list[list[str]] = []
    for r in rows_data:
        table_rows.append(
            [
                r.get("category", "")[:14],
                (r.get("query_text", "")[:42] + "…") if len(r.get("query_text", "")) > 42 else r.get("query_text", ""),
                r.get("amem_rank", ""),
                r.get("tmem_rank", ""),
                r.get("movement_vs_amem", ""),
            ]
        )
    if table_rows:
        _table_slide(
            prs,
            "Results: ground-truth rank (lower is better)",
            ["Category", "Query (trunc.)", "A-MEM", "T-MEM", "vs A-MEM"],
            table_rows,
        )

    _bullet_slide(
        prs,
        "Aggregate (this run)",
        [
            "T-MEM improved correct-answer best rank on 6 / 7 queries; 1 tie; 0 worse.",
            "Mean rank delta (A-MEM rank − T-MEM rank): ≈ 2.3 positions.",
            "Per-category @1/@3 accuracy: see chart + summary_stats.txt in presentation folder.",
            "Re-run with real MiniLM (omit TMEM_ALLOW_HASH_ENCODER) for paper-aligned embeddings.",
        ],
    )

    _image_slide(prs, "Retrieval accuracy by category", HERE / "retrieval_accuracy.png")
    _image_slide(prs, "Score breakdown: “Where does Alice currently work?”", HERE / "score_breakdown.png")
    _image_slide(prs, "Decay, reinforcement, and link scenarios (illustrative)", HERE / "decay_curves.png", max_h=4.9)

    _bullet_slide(
        prs,
        "Takeaways & limitations",
        [
            "T-MEM is a controlled extension: same note pipeline, modified ranking only.",
            "Decay uses month-equivalent age so λ stays meaningful over long corpora — document in thesis.",
            "Global r̂ can favor high-access topics; tune β / access simulation for your data.",
            "Baseline retrieval is full-memory cosine (Eq. 9), not Chroma neighbor expansion from A-mem-sys.",
        ],
    )

    out = HERE / "T_MEM_Deck.pptx"
    prs.save(str(out))
    print(f"Wrote {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
